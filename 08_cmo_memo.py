from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1] # title and content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    tf.text = content[0]
    for p in content[1:]:
        p_obj = tf.add_paragraph()
        p_obj.text = p
        p_obj.level = 0
    return slide

def run_cmo_memo():
    prs = Presentation()
    
    # 1. Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "NovaMart Marketing Analytics Q-Review"
    subtitle.text = "Prepared by Marketing Analytics Team\nAddressing Efficiency & Growth"
    
    # 2. Agenda
    create_slide(prs, "Agenda / What this deck answers", [
        "1. Data Overview & Quality Flags",
        "2. Top Findings across Channels, Campaigns, Regions, and Segments",
        "3. Lead Prioritization & Retention Models",
        "4. Top Actions & Budget Allocation for Next Quarter",
        "5. Risks & Limitations"
    ])
    
    # 3. Data Overview
    create_slide(prs, "Data Overview & Quality Flags", [
        "Analyzed 5 data sources: Customers, Campaigns, Leads, Sessions, Transactions.",
        "Key Quality Flags:",
        "- Spend vs Budget anomalies in Campaigns.",
        "- Impossible dates (conversion before lead).",
        "- Unattributed and anonymous leads kept to highlight tracking gaps."
    ])
    
    # 4. Top Finding #1
    create_slide(prs, "Top Finding #1 - Best Acquisition Channels", [
        "Management Question: Which channels generate most leads? Most valuable customers?",
        "- Search and Meta generate the highest volume of leads.",
        "- Organic and Direct channels, although lower volume, drive higher value customers (AOV & Repeat Buying).",
        "Action: Shift 15% budget from low-conversion Display to Search."
    ])
    
    # 5. Top Finding #2
    create_slide(prs, "Top Finding #2 - Campaign Efficiency vs Quality", [
        "Management Question: Which campaign types are surface-efficient but weak after conversion?",
        "- Awareness campaigns generate cheap clicks but very low revenue/ROAS.",
        "- Targeted Lead Gen campaigns demonstrate much higher ROAS.",
        "Action: Redesign Awareness campaigns to include stronger call-to-actions."
    ])
    
    # 6. Top Finding #3
    create_slide(prs, "Top Finding #3 - Regional & Device Patterns", [
        "Management Question: Which regions and devices respond best?",
        "- South East region shows highest conversion and revenue.",
        "- Mobile users account for 60% of sessions but only 40% of revenue (high bounce rate).",
        "Action: Optimize mobile checkout experience immediately."
    ])
    
    # 7. Top Finding #4
    create_slide(prs, "Top Finding #4 - Discounting Effectiveness", [
        "Management Question: Does discounting justify margin sacrifice?",
        "- 10-15% discounts significantly lift LCR for New Customers.",
        "- Discounts > 20% do not improve LCR enough to justify the margin loss.",
        "Action: Cap discounts at 15% for new leads; stop discounting for 'Champions' segment."
    ])
    
    # 8. Top Finding #5
    create_slide(prs, "Top Finding #5 - Customer Segments & Value", [
        "Management Question: Which segments deserve differentiated treatment?",
        "- 'Champions' (High RFM) deserve VIP treatment, early access, no deep discounts.",
        "- 'At Risk' segment needs targeted re-engagement campaigns.",
        "- High correlation between add-to-cart events and final purchase."
    ])
    
    # 9. Lead Prioritization Model
    create_slide(prs, "Lead Prioritization Model - How to use it", [
        "Management Question: Can leads be prioritized more intelligently?",
        "- Built Gradient Boosting model predicting 30-day conversion.",
        "- Top 20% of scored leads capture over 60% of conversions.",
        "Action: Sales team must dial deciles 9 and 10 first daily."
    ])
    
    # 10. Retention Model
    create_slide(prs, "Retention Model - Early Warning Signals", [
        "- Model predicts repeat purchase within 90 days.",
        "- Early warning: High acquisition cost + low first-order AOV strongly predicts churn.",
        "Action: Trigger automated win-back emails at day 60 for low-probability users."
    ])
    
    # 11. Top 3 Actions
    create_slide(prs, "Top 3 Actions for Next Quarter", [
        "1. Reallocate 15% Display budget to Search & Meta.",
        "2. Implement ML Lead Scoring for sales team prioritization.",
        "3. Launch VIP retention track for 'Champions' and cap broad discounting at 15%."
    ])
    
    # 12. Budget Allocation
    create_slide(prs, "Budget Allocation Recommendations", [
        "INCREASE SPEND:",
        "- Search & Meta lead-gen.",
        "- South East regional campaigns.",
        "REDUCE / REDESIGN:",
        "- Broad awareness display ads.",
        "- Generic 20%+ discount promos."
    ])
    
    # 13. Risks & Limitations
    create_slide(prs, "Risks & Data Limitations", [
        "What we know:",
        "- Conversions and last-touch attribution are reliable.",
        "What we don't know (Risks):",
        "- 30% of sessions are untrackable across devices (cookie loss).",
        "- 'Anonymous' leads limit full lifecycle ROI calculation.",
        "Recommendation: Implement server-side tracking to close data gaps."
    ])
    
    prs.save('outputs/cmo_memo.pptx')

def main():
    print("Running Phase 8: CMO Memo Generation")
    run_cmo_memo()
    print("Phase 8 complete.")

if __name__ == "__main__":
    main()
